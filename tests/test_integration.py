"""Integracao real: PPTX -> LibreOffice -> PDF -> cortes.

Marcado como `slow` porque sobe o LibreOffice. Rodar com `-m "not slow"` para pular.
"""

from __future__ import annotations

import pytest

from slidecut import cli, convert

pptx = pytest.importorskip("pptx")

pytestmark = pytest.mark.skipif(
    convert.find_soffice() is None, reason="LibreOffice nao instalado"
)

ORANGE = pptx.dml.color.RGBColor(0xB0, 0x6E, 0x03)
SLIDES = [
    ("Conceito", True),
    ("conteudo um", False),
    ("conteudo dois", False),
    ("Fontes", True),
    ("conteudo tres", False),
]


def _build_pptx(path):
    from pptx import Presentation
    from pptx.util import Pt

    deck = Presentation()
    blank = deck.slide_layouts[6]
    for text, is_divider in SLIDES:
        slide = deck.slides.add_slide(blank)
        if is_divider:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = ORANGE
        box = slide.shapes.add_textbox(Pt(80), Pt(160), Pt(560), Pt(100))
        frame = box.text_frame
        frame.text = text
        frame.paragraphs[0].runs[0].font.size = Pt(40)
    deck.save(str(path))
    return path


@pytest.mark.slow
def test_pptx_is_converted_and_split_by_divider_slides(tmp_path, capsys):
    source = _build_pptx(tmp_path / "aula.pptx")
    outdir = tmp_path / "out"

    assert cli.main([str(source), "-o", str(outdir)]) == 0

    produced = sorted(p.name for p in outdir.glob("*.pdf"))
    assert produced == ["01 - Conceito.pdf", "02 - Fontes.pdf"]
    assert "Cor divisora: #B0" in capsys.readouterr().out
